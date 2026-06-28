"""
Fill the mandatory H2S/Redrob Google-Slides template with our approach content.
Reads deck/h2s_template.pptx (the official template, exported from Google Slides),
writes deck/Redrob_H2S_Deck.pptx. The user then opens it in Google Slides /
PowerPoint, eyeballs it, and exports a PDF (<=5 MB) for the portal.

    python scripts/fill_h2s_deck.py
"""
from __future__ import annotations
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

TPL = "deck/h2s_template.pptx"
OUT = "deck/Redrob_H2S_Deck.pptx"

INK = RGBColor(0x20, 0x27, 0x29)      # template's dark charcoal body colour
BODY_FONT = "Manrope SemiBold"
LEAD_FONT = "Manrope ExtraBold"
BODY_PT = 12

# --- per-slide content: list of (bold_lead, body) paragraphs -----------------
CONTENT = {
    2: [  # Solution Overview
        ("Our solution.  ", "A hybrid, CPU-only ranker that reads each profile the way a recruiter would — scoring fit from evidence, not keyword counts."),
        ("Scoring.  ", "score = (0.70 · structured JD-fit + 0.30 · contrastive semantic fit) × behavioral availability × honeypot gate."),
        ("What makes it different.  ", "It is built to beat the dataset's traps — keyword stuffers, plain-language strong candidates, behavioral twins, and ~80 honeypots."),
        ("No LLM at inference.  ", "Deterministic, inspectable and reproducible on CPU in ~3 minutes; it distrusts the skills tag-cloud and reads career evidence instead."),
    ],
    3: [  # JD Understanding & Candidate Evaluation
        ("Requirements read from the JD.  ", "6–8 yrs applied ML at product (not services) companies; shipped search / ranking / retrieval / recsys to real users; embeddings + vector/hybrid search with rigorous NDCG/MAP-style evaluation; strong NLP/IR; near Pune–Noida; actually available."),
        ("Signals that matter most.  ", "Title coherence, demonstrated production career evidence, trust-weighted skills, and availability."),
        ("Beyond keyword matching.  ", "A contrastive embedding surfaces people who describe building these systems without naming the buzzwords, while explicit anti-patterns — research-only, LangChain-on-OpenAI wrappers, title-chasers — are penalised."),
    ],
    4: [  # Ranking Methodology
        ("Retrieve, score, rank.  ", "Stream all 100K profiles, compute a per-candidate score in one pass, then select and order the top-100 (ties broken by candidate_id)."),
        ("Structured fit (0.70).  ", "Title taxonomy, career-evidence patterns read from descriptions, skill-trust = endorsements × duration × proficiency × assessment, plus experience band, domain and location."),
        ("Semantic fit (0.30).  ", "cos(candidate, JD-ideal) − cos(candidate, JD-anti-pattern) using model2vec static embeddings, with a pure scikit-learn LSA fallback."),
        ("Combining signals.  ", "Weighted blend × behavioral availability modifier × honeypot gate."),
    ],
    5: [  # Explainability & Data Validation
        ("Every pick is explained.  ", "Each candidate ships fact-grounded reasoning plus interpretable sub-scores (title fit, career evidence, skill trust, availability)."),
        ("No hallucinations.  ", "Reasoning is generated deterministically from the candidate's own parsed fields — no LLM, nothing invented — and honest concerns are surfaced."),
        ("Suspicious / low-quality profiles.  ", "A honeypot gate flags internal impossibilities (tenure > time elapsed; ≥3 'expert' skills with 0 months of use); skill-trust defeats stuffing; dormant or unresponsive profiles are down-weighted."),
    ],
    6: [  # End-to-End Workflow
        ("", "JD encoded as taxonomies, anchors and weights  →  stream candidates.jsonl  →  per candidate: build text + evidence blob  →  structured score + semantic score + behavioral modifier + honeypot check  →  blend into one score  →  select & order top-100  →  emit submission.csv with grounded reasoning."),
        ("One command, offline.  ", "python rank.py --candidates candidates.jsonl --out submission.csv   ·   CPU-only   ·   ~3 min."),
    ],
    7: [  # System Architecture (no body box in template -> we add one)
        ("Inputs.  ", "candidates.jsonl (100K profiles) + the JD encoded in config (taxonomies, evidence patterns, anchors, weights)."),
        ("Layer 1 — Structured fit (0.70).  ", "title · career evidence · skill trust · experience · domain · location."),
        ("Layer 2 — Semantic fit (0.30).  ", "contrastive static embeddings (model2vec, scikit-learn LSA fallback)."),
        ("Layer 3 — Behavioral modifier.  ", "availability from the 23 redrob signals."),
        ("Layer 4 — Honeypot gate.  ", "removes internally-impossible profiles from the top-100."),
        ("Output.  ", "ranked top-100 + grounded reasoning. Two-pass: scalar scoring over 100K, then full facts for the top 100."),
    ],
    8: [  # Results & Performance
        ("Top-100 quality.  ", "95% India-based  ·  85% open-to-work  ·  6.4 yrs mean experience  ·  0 honeypots (cap is 10%). Scores span 0.99 → 0.88."),
        ("Beats the keyword trap.  ", "The shipped sample ranks an HR Manager #1; our top picks are Recommendation-Systems, Search and Applied-ML Engineers. Passes the official validator."),
        ("Runtime & compute.  ", "CPU-only, no GPU, no network at ranking; ~3.2 min for 100K profiles (< 5-min budget); ≤16 GB via streaming. Embeddings precomputed offline."),
    ],
    9: [  # Technologies Used
        ("Core.  ", "Python 3.12  ·  NumPy  ·  scikit-learn (TF-IDF + TruncatedSVD LSA)  ·  model2vec static embeddings  ·  orjson  ·  PyYAML."),
        ("Why these.  ", "The CPU-only / offline constraint rules out transformers — static embeddings give real semantics at hashing speed with no torch or GPU; scikit-learn keeps it deterministic and reproducible; orjson streams 100K profiles fast."),
        ("Delivery.  ", "Streamlit (hosted sandbox)  ·  reportlab  ·  Docker for one-command reproduction."),
    ],
    10: [  # Submission Assets
        ("GitHub (public).  ", "https://github.com/PureGenius369/redrob-ranker"),
        ("Live sandbox (HF Spaces).  ", "https://huggingface.co/spaces/MannSutariya/redrob-ranker"),
        ("Ranked output.  ", "submission.xlsx — top-100 with candidate_id, rank, score and grounded reasoning."),
        ("Reproduce.  ", "python rank.py --candidates candidates.jsonl --out submission.csv"),
    ],
}

COVER = {  # slide 1 label -> value appended after the colon
    "Team Name :": "Mann Sutariya",
    "Team Leader Name :": "Mann Sutariya",
    "Problem Statement :": "Intelligent Candidate Discovery & Ranking (Data & AI Challenge)",
}


def style_run(r, *, lead=False, size=BODY_PT):
    r.font.name = LEAD_FONT if lead else BODY_FONT
    r.font.size = Pt(size)
    r.font.color.rgb = INK
    r.font.bold = True if lead else None


def fill_body(tf, paras):
    tf.word_wrap = True
    tf.clear()
    for idx, (lead, body) in enumerate(paras):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        p.line_spacing = 1.08
        if lead:
            rl = p.add_run(); rl.text = lead; style_run(rl, lead=True)
        rb = p.add_run(); rb.text = body; style_run(rb, lead=False)


def main():
    prs = Presentation(TPL)
    slides = list(prs.slides)

    # ---- slide 1: cover ----
    for sh in slides[0].shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t in COVER:
                para = sh.text_frame.paragraphs[0]
                # find the label run's font to match the appended value
                base = next((r for r in para.runs if r.text.strip()), None)
                run = para.add_run()
                run.text = " " + COVER[t]
                if base is not None:
                    run.font.name = base.font.name
                    if base.font.size:
                        run.font.size = base.font.size
                    run.font.color.rgb = INK

    # ---- content slides ----
    for n, paras in CONTENT.items():
        slide = slides[n - 1]
        # body box = the text box that sits below the title (top > 1.2")
        body_box = None
        for sh in slide.shapes:
            if sh.has_text_frame and sh.top is not None and sh.top > Inches(1.2):
                body_box = sh
                break
        if body_box is None:  # slide 7 has only a title -> add a body text box
            body_box = slide.shapes.add_textbox(Inches(0.41), Inches(1.4),
                                                Inches(9.32), Inches(3.5))
        fill_body(body_box.text_frame, paras)

    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
